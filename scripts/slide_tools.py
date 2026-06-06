#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
yoko-slides 後処理ツール（やさしい学術教科書スライド）

HiggsField(nano_banana_pro)で生成した16:9スライド群を、納品形に整える再利用CLI。
Y-007 で手作業した処理を部品化したもの。

サブコマンド:
  whitefix    純白の箱/帯をクリーム地に置換（ピンク枠内の白文字は保護）
  telopsafe   84%縮小＋同色余白で下部15%のテロップ余白を確保
  stamp       金の中抜き番号(01等)を左上に合成（Didot, 全スライド同一デザイン）
  pptx        16:9フルブリードPPTXを生成
  finish      whitefix→telopsafe→pptx を一括（番号は別途stampしてから流す）

使い方例:
  python3 slide_tools.py whitefix raw/ clean/
  python3 slide_tools.py telopsafe clean/ final/
  python3 slide_tools.py stamp raw/slide_06.png raw/slide_06.png --num 01
  python3 slide_tools.py pptx final/ deck.pptx
  python3 slide_tools.py finish raw/ final/ deck.pptx
"""
import argparse, glob, os, sys
from collections import Counter
from PIL import Image, ImageDraw, ImageFont
import numpy as np

# ---- 共通定数（マスター規格と一致させること） ----
GOLD = (201, 161, 74)
PINK = (232, 74, 138)
INDIGO = (46, 58, 89)
FONT_NUM_CANDIDATES = [
    "/System/Library/Fonts/Supplemental/Didot.ttc",
    "/System/Library/Fonts/Supplemental/Bodoni 72.ttc",
    "/System/Library/Fonts/Supplemental/Georgia.ttf",
]
FONT_JP_CANDIDATES = [
    "/System/Library/Fonts/ヒラギノ角ゴシック W6.ttc",
    "/System/Library/Fonts/ヒラギノ角ゴシック W7.ttc",
    "/System/Library/Fonts/Hiragino Sans GB.ttc",
]

def _font(cands, size):
    for p in cands:
        if os.path.exists(p):
            return ImageFont.truetype(p, size)
    raise FileNotFoundError("font not found: " + ", ".join(cands))

def bg_color(im):
    """上端1%の最頻色＝背景クリーム色を推定"""
    im = im.convert("RGB")
    W, H = im.size; px = im.load(); c = Counter()
    for y in range(0, max(2, int(H * 0.01))):
        for x in range(0, W, 5):
            c[px[x, y]] += 1
    return c.most_common(1)[0][0]

def _iter(indir):
    if os.path.isfile(indir):
        return [indir]
    fs = sorted(glob.glob(os.path.join(indir, "*.png")))
    return [f for f in fs if not os.path.basename(f).startswith("_")]

# ---------- whitefix ----------
def remove_white(im, dilate=6):
    """純白(min>=248 かつ無彩色)をbgに。ピンク領域を膨張保護して白文字は残す。"""
    im = im.convert("RGB")
    bg = bg_color(im)
    arr = np.array(im).astype(np.int16)
    mn = arr.min(2); mx = arr.max(2)
    white = (mn >= 248) & ((mx - mn) < 12)
    R, G, B = arr[:, :, 0], arr[:, :, 1], arr[:, :, 2]
    pink = (R > 175) & ((R - G) > 45) & ((R - B) > 25)
    prot = pink.copy()
    for _ in range(dilate):
        p = prot.copy()
        p[1:, :] |= prot[:-1, :]; p[:-1, :] |= prot[1:, :]
        p[:, 1:] |= prot[:, :-1]; p[:, :-1] |= prot[:, 1:]
        prot = p
    out = np.array(im)
    out[white & (~prot)] = bg
    return Image.fromarray(out)

# ---------- telopsafe ----------
def telop_safe(im, scale=0.84, top_frac=0.015):
    """中身を縮小して上寄せ、下部15%を同色余白に。黒帯にならず自然に余白確保。"""
    im = im.convert("RGB"); W, H = im.size; bg = bg_color(im)
    canvas = Image.new("RGB", (W, H), bg)
    sw, sh = int(W * scale), int(H * scale)
    sc = im.resize((sw, sh), Image.LANCZOS)
    x = (W - sw) // 2; y = int(H * top_frac)
    if y + sh > int(H * 0.85):
        y = max(0, int(H * 0.85) - sh)
    canvas.paste(sc, (x, y))
    return canvas

# ---------- stamp number ----------
def stamp_number(im, num, fs_frac=0.20, x_frac=0.045, y_frac=0.05):
    """左上に金の中抜き番号。全スライドで同一の見た目になる。"""
    im = im.convert("RGB"); W, H = im.size; bg = bg_color(im)
    d = ImageDraw.Draw(im)
    fs = int(H * fs_frac); f = _font(FONT_NUM_CANDIDATES, fs)
    d.text((int(W * x_frac), int(H * y_frac)), num, font=f,
           fill=bg, stroke_width=max(3, fs // 40), stroke_fill=GOLD)
    return im

# ---------- pptx ----------
def build_pptx(indir, out):
    from pptx import Presentation
    from pptx.util import Inches
    files = _iter(indir)
    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for f in files:
        s = prs.slides.add_slide(blank)
        s.shapes.add_picture(f, 0, 0, width=prs.slide_width, height=prs.slide_height)
    prs.save(out)
    return len(files)

# ---------- audit ----------
def audit_bottom(im):
    """下部15%に内容が食い込んでいないか（content最下端%）"""
    im = im.convert("RGB"); W, H = im.size; px = im.load(); bg = px[W // 2, 4]
    def diff(c): return abs(c[0]-bg[0])+abs(c[1]-bg[1])+abs(c[2]-bg[2])
    for y in range(H - 1, -1, -1):
        cnt = sum(1 for x in range(0, W, 3) if diff(px[x, y]) > 40)
        if cnt > W * 0.02:
            return 100 * y / H
    return 0.0

def main():
    ap = argparse.ArgumentParser()
    sub = ap.add_subparsers(dest="cmd", required=True)
    for name in ("whitefix", "telopsafe"):
        s = sub.add_parser(name); s.add_argument("indir"); s.add_argument("outdir")
    sub._name_parser_map["telopsafe"].add_argument("--scale", type=float, default=0.84)
    sp = sub.add_parser("stamp"); sp.add_argument("infile"); sp.add_argument("outfile"); sp.add_argument("--num", required=True)
    pp = sub.add_parser("pptx"); pp.add_argument("indir"); pp.add_argument("out")
    fn = sub.add_parser("finish"); fn.add_argument("indir"); fn.add_argument("outdir"); fn.add_argument("out"); fn.add_argument("--scale", type=float, default=0.84)
    au = sub.add_parser("audit"); au.add_argument("indir")
    a = ap.parse_args()

    if a.cmd == "whitefix":
        os.makedirs(a.outdir, exist_ok=True)
        for f in _iter(a.indir):
            remove_white(Image.open(f)).save(os.path.join(a.outdir, os.path.basename(f)))
            print("whitefix:", os.path.basename(f))
    elif a.cmd == "telopsafe":
        os.makedirs(a.outdir, exist_ok=True)
        for f in _iter(a.indir):
            telop_safe(Image.open(f), scale=a.scale).save(os.path.join(a.outdir, os.path.basename(f)))
            print("telopsafe:", os.path.basename(f))
    elif a.cmd == "stamp":
        stamp_number(Image.open(a.infile), a.num).save(a.outfile)
        print("stamp", a.num, "->", a.outfile)
    elif a.cmd == "pptx":
        n = build_pptx(a.indir, a.out); print(f"pptx: {n}枚 -> {a.out}")
    elif a.cmd == "finish":
        tmp = os.path.join(a.outdir, "_clean"); os.makedirs(tmp, exist_ok=True)
        os.makedirs(a.outdir, exist_ok=True)
        for f in _iter(a.indir):
            im = remove_white(Image.open(f))
            im = telop_safe(im, scale=a.scale)
            im.save(os.path.join(a.outdir, os.path.basename(f)))
        n = build_pptx(a.outdir, a.out); print(f"finish: {n}枚 -> {a.out}")
    elif a.cmd == "audit":
        bad = 0
        for f in _iter(a.indir):
            p = audit_bottom(Image.open(f))
            flag = "OK" if p <= 85.6 else "★余白不足"
            if p > 85.6: bad += 1
            print(f"{os.path.basename(f):40} 内容最下端 {p:5.1f}%  {flag}")
        print(f"余白不足: {bad}枚")

if __name__ == "__main__":
    main()
